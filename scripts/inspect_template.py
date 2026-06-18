"""Inspeciona o template Status-Report.pptx local para mapear placeholders."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation


def main(pptx_path: str) -> int:
    p = Presentation(pptx_path)
    print(f"Arquivo: {pptx_path}")
    print(f"Total de slides: {len(p.slides)}")
    for i, slide in enumerate(p.slides, 1):
        print(f"\n=== Slide {i} ===")
        for shape in slide.shapes:
            print(f"- shape: {shape.shape_type} | name={shape.name}")
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        print(f"    texto: {txt!r}")
            if getattr(shape, "has_table", False):
                print("    [TABELA]")
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    print("     | ".join(cells))
    return 0


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else r"C:\\Users\\pedro\\Downloads\\Status-Report.pptx"
    sys.exit(main(path))
